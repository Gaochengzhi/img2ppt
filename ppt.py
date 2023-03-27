import collections
import collections.abc
import os
from pptx import Presentation
from pptx.util import Inches
import re
import glob

# config your params
image_folder = "."
w = 16
h = 12
file_name="output.pptx"

presentation = Presentation()
# Set slide width and height (optional)
presentation.slide_width = Inches(w)
presentation.slide_height = Inches(h)

image_files = (
    glob.glob(os.path.join(image_folder, "*.png"))
    + glob.glob(os.path.join(image_folder, "*.jpg"))
    + glob.glob(os.path.join(image_folder, "*.jpeg"))
)


def extract_number(filename):
    number = re.search(r"\d+", os.path.basename(filename))
    return int(number.group()) if number else float("inf")


image_files.sort(key=extract_number)


# Add slides with images

for image_file in image_files:
    # Create a new slide with a blank layout
    # slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    # Check if the image file exists
    if os.path.exists(image_file):
        # Create a new slide with a blank layout
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])

        # Add the image to the slide
        left, top, width, height = Inches(0), Inches(0), Inches(w), Inches(h)
        slide.shapes.add_picture(image_file, left, top, width, height)
    else:
        print(f"Image {image_file} not found. Skipping...")

# Save the presentation
presentation.save(file_name)
print(f"PowerPoint presentation created: {file_name}")

